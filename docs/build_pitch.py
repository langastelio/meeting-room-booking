# -*- coding: utf-8 -*-
"""Generates a professional pitch deck for the PBB Leadership presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (Standard Bank corporate blues) ----
NAVY   = RGBColor(0x00, 0x1A, 0x4B)   # deep navy
BLUE   = RGBColor(0x00, 0x33, 0xA0)   # SB blue
ACCENT = RGBColor(0x00, 0xA6, 0xE2)   # bright cyan accent
LIGHT  = RGBColor(0xF2, 0xF5, 0xFB)   # near-white panel
GREY   = RGBColor(0x5A, 0x63, 0x72)   # body grey
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines = list of (text, size, bold, color, space_after)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, sa) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sa)
        r = p.add_run()
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def bullet(slide, x, y, w, h, items, size=16, color=GREY, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "–  " + head
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = NAVY
        if body:
            r2 = p.add_run(); r2.text = body
            r2.font.name = FONT; r2.font.size = Pt(size); r2.font.bold = False; r2.font.color.rgb = color
    return tb


def footer(slide, page):
    txt(slide, Inches(0.55), Inches(7.02), Inches(9), Inches(0.35),
        [("Meeting Room Booking  ·  PPB Moçambique", 10, False, GREY, 0)])
    txt(slide, Inches(11.8), Inches(7.02), Inches(1.0), Inches(0.35),
        [(str(page), 10, True, BLUE, 0)], align=PP_ALIGN.RIGHT)


def accent_bar(slide):
    rect(slide, 0, 0, SW, Inches(0.12), ACCENT)


# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, SW, Inches(0.18), ACCENT)
# side accent block
rect(s, 0, Inches(2.4), Inches(0.35), Inches(2.4), BLUE)
txt(s, Inches(0.9), Inches(2.35), Inches(11), Inches(0.6),
    [("▦  MEETING ROOM BOOKING", 20, True, ACCENT, 0)])
txt(s, Inches(0.85), Inches(2.95), Inches(11.6), Inches(2.0),
    [("Reserva de Salas, do zero à produção", 40, True, WHITE, 6),
     ("Construído com Inteligência Artificial em cada passo", 22, False, RGBColor(0xC8,0xD6,0xF0), 0)])
rect(s, Inches(0.9), Inches(5.15), Inches(3.2), Inches(0.045), ACCENT)
txt(s, Inches(0.9), Inches(5.35), Inches(11), Inches(0.8),
    [("Apresentação à Liderança do PBB", 15, False, RGBColor(0xC8,0xD6,0xF0), 3),
     ("Standard Bank  ·  PPB Moçambique", 13, False, GREY, 0)])

# =====================================================================
# SLIDE 2 — Problem -> Solution
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
accent_bar(s)
txt(s, Inches(0.55), Inches(0.45), Inches(12), Inches(0.9),
    [("O problema → a solução", 30, True, NAVY, 0)])
rect(s, Inches(0.55), Inches(1.25), Inches(1.4), Inches(0.05), ACCENT)

# Before card
rect(s, Inches(0.55), Inches(1.7), Inches(5.8), Inches(4.6), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.9), Inches(1.95), Inches(5.2), Inches(0.6),
    [("ANTES", 16, True, GREY, 0)])
bullet(s, Inches(0.9), Inches(2.65), Inches(5.2), Inches(3.4), [
    ("Gestão por email ", "e mensagens dispersas"),
    ("Conflitos de reserva ", "constantes na mesma sala"),
    ("Sem visibilidade ", "de ocupação ou histórico"),
    ("Sem relatórios ", "para a gestão"),
], size=15, gap=13)

# After card
rect(s, Inches(6.95), Inches(1.7), Inches(5.8), Inches(4.6), BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.3), Inches(1.95), Inches(5.2), Inches(0.6),
    [("DEPOIS", 16, True, ACCENT, 0)])
tb = s.shapes.add_textbox(Inches(7.3), Inches(2.65), Inches(5.2), Inches(3.4))
tf = tb.text_frame; tf.word_wrap = True
after = [
    ("App web completa ", "de reserva de salas"),
    ("Calendário ", "dia / semana / mês"),
    ("Deteta conflitos ", "automaticamente"),
    ("Painel + relatórios ", "e impressão de escalas"),
    ("Bilingue ", "Português e Inglês"),
]
for i,(h,b) in enumerate(after):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after = Pt(11)
    r=p.add_run(); r.text="✓  "+h; r.font.name=FONT; r.font.size=Pt(15); r.font.bold=True; r.font.color.rgb=WHITE
    r2=p.add_run(); r2.text=b; r2.font.name=FONT; r2.font.size=Pt(15); r2.font.color.rgb=RGBColor(0xD5,0xE2,0xF7)
footer(s, 2)

# =====================================================================
# SLIDE 3 — Tech stack & why (AI-guided)
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
accent_bar(s)
txt(s, Inches(0.55), Inches(0.45), Inches(12), Inches(0.9),
    [("A tecnologia — e o porquê", 30, True, NAVY, 0)])
rect(s, Inches(0.55), Inches(1.25), Inches(1.4), Inches(0.05), ACCENT)
txt(s, Inches(0.55), Inches(1.45), Inches(12), Inches(0.6),
    [("Escolha da stack guiada por AI: rápida de entregar, sem licenças, fácil de manter.", 15, False, GREY, 0)])

cards = [
    ("HTML · CSS · JS", "Front-end puro. Sem frameworks pesadas, corre em qualquer navegador."),
    ("Supabase", "Base de dados na cloud, partilhada e ao vivo. Zero servidores a gerir."),
    ("Custo ~ zero", "Assente em tiers gratuitos. Sem custos de licenciamento."),
]
cx = Inches(0.55); cw = Inches(3.95); gap = Inches(0.25)
for i,(h,b) in enumerate(cards):
    x = Emu(int(cx) + i*(int(cw)+int(gap)))
    rect(s, x, Inches(2.35), cw, Inches(2.9), LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, Inches(2.35), cw, Inches(0.12), ACCENT, MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Emu(int(x)+Inches(0.3)), Inches(2.7), Emu(int(cw)-Inches(0.6)), Inches(0.7),
        [(h, 19, True, BLUE, 0)])
    txt(s, Emu(int(x)+Inches(0.3)), Inches(3.5), Emu(int(cw)-Inches(0.6)), Inches(1.5),
        [(b, 14, False, GREY, 0)])

rect(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.0), NAVY, MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.9), Inches(5.72), Inches(11.6), Inches(0.8),
    [("Eu decidi o QUÊ e validei a qualidade. A AI comparou alternativas e acelerou o COMO.",
      16, True, WHITE, 0)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# =====================================================================
# SLIDE 4 — AI in every step
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
accent_bar(s)
txt(s, Inches(0.55), Inches(0.45), Inches(12), Inches(0.9),
    [("AI em cada passo", 30, True, NAVY, 0)])
rect(s, Inches(0.55), Inches(1.25), Inches(1.4), Inches(0.05), ACCENT)

steps = [
    ("1", "DESENHAR", "Modelo de dados e regras de negócio — ex.: impedir reservas sobrepostas."),
    ("2", "PROGRAMAR", "Cada página: calendário, painel de estatísticas, gestão de utilizadores."),
    ("3", "TRADUZIR", "Aplicação inteira localizada para Português e Inglês."),
    ("4", "EVOLUIR", "Funcionalidades novas em minutos — ex.: imprimir escala por sala."),
]
gy = Inches(1.75); ch = Inches(1.2); vg = Inches(0.12)
for i,(n,h,b) in enumerate(steps):
    y = Emu(int(gy) + i*(int(ch)+int(vg)))
    rect(s, Inches(0.55), y, Inches(12.2), ch, LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.55), y, Inches(1.2), ch, BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.55), y, Inches(1.2), ch,
        [(n, 34, True, WHITE, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(2.0), y, Inches(3.0), ch,
        [(h, 19, True, NAVY, 0)], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.0), y, Inches(7.5), ch,
        [(b, 15, False, GREY, 0)], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# =====================================================================
# SLIDE 5 — Why only possible with AI + vision
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, NAVY)
rect(s, 0, 0, SW, Inches(0.18), ACCENT)
txt(s, Inches(0.85), Inches(0.7), Inches(11.6), Inches(0.9),
    [("Porque só foi possível com AI", 30, True, WHITE, 0)])
rect(s, Inches(0.9), Inches(1.5), Inches(1.4), Inches(0.05), ACCENT)

pts = [
    ("Um projeto destes exigiria uma equipa e semanas.", "Entregue por uma pessoa, em dias."),
    ("A AI foi programador, revisor e tradutor.", "Disponível a qualquer hora."),
    ("Eu mantive-me o arquiteto.", "Decidi o quê, validei a qualidade."),
]
tb = s.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
for i,(h,b) in enumerate(pts):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after = Pt(14)
    r=p.add_run(); r.text="▸  "+h+"  "; r.font.name=FONT; r.font.size=Pt(19); r.font.bold=True; r.font.color.rgb=ACCENT
    r2=p.add_run(); r2.text=b; r2.font.name=FONT; r2.font.size=Pt(19); r2.font.color.rgb=RGBColor(0xD5,0xE2,0xF7)

rect(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.6), BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.3), Inches(5.05), Inches(10.8), Inches(1.3),
    [("Não é substituir pessoas — é torná-las 10x mais rápidas.", 22, True, WHITE, 6),
     ("Se um projeto assim se entrega tão depressa, imaginem as equipas do PBB com estas ferramentas.",
      15, False, RGBColor(0xD5,0xE2,0xF7), 0)], anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.5),
    [("Obrigado.", 16, True, ACCENT, 0)])

out = r"docs\Apresentacao-Lideranca-PBB.pptx"
prs.save(out)
print("Saved:", out, "-", len(prs.slides._sldIdLst), "slides")
